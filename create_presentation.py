"""Generate Sprint 1 Review presentation for CI Testing project."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
DARK_CARD = RGBColor(0x25, 0x25, 0x3F)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)
CODE_FG = RGBColor(0xE6, 0xED, 0xF3)
KW_COLOR = RGBColor(0xFF, 0x7B, 0x72)    # keywords (red-ish)
STR_COLOR = RGBColor(0xA5, 0xD6, 0xFF)   # strings/types (blue)
CMT_COLOR = RGBColor(0x8B, 0x94, 0x9E)   # comments (gray)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = alignment
    return tf


def add_bullets(slide, items, left, top, width, height, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.text = item
        p.font.color.rgb = color
        p.font.size = Pt(font_size)
        p.font.name = "Segoe UI"
    return tf


def add_code_block(slide, lines, left, top, width, height, font_size=13,
                   rust_keywords=None, py_keywords=None):
    """Add a styled code block with dark background and monospace font.

    lines can be strings or tuples of (text, color) for manual coloring.
    If rust_keywords or py_keywords are provided, basic syntax highlighting is applied.
    """
    import re
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    # Add border accent
    shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    shape.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(left + Inches(0.4), top + Inches(0.35),
                                      width - Inches(0.8), height - Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True

    rust_kw = {"trait", "fn", "pub", "impl", "struct", "let", "mut", "use", "self", "&self",
               "mod", "enum", "for", "if", "else", "return", "match", "Ok", "Err", "true", "false"}
    py_kw = {"from", "import", "def", "class", "return", "True", "False", "None", "as", "if",
             "else", "for", "in", "with", "not", "and", "or", "is"}

    keywords = set()
    if rust_keywords:
        keywords = rust_kw
    if py_keywords:
        keywords = py_kw

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(0)

        if isinstance(line, tuple):
            run = p.add_run()
            run.text = line[0]
            run.font.color.rgb = line[1]
            run.font.size = Pt(font_size)
            run.font.name = "Cascadia Code"
            continue

        if not line.strip():
            run = p.add_run()
            run.text = " "
            run.font.size = Pt(font_size)
            run.font.name = "Cascadia Code"
            continue

        # Check if it's a comment line (but not Rust attributes like #[...])
        stripped = line.lstrip()
        if stripped.startswith("//") or (stripped.startswith("#") and not stripped.startswith("#[")):
            run = p.add_run()
            run.text = line
            run.font.color.rgb = CMT_COLOR
            run.font.size = Pt(font_size)
            run.font.name = "Cascadia Code"
            continue

        # Tokenize and colorize
        tokens = re.split(r'(\s+|[().,:{}\[\]=<>])', line)
        for token in tokens:
            if not token:
                continue
            run = p.add_run()
            run.text = token
            run.font.size = Pt(font_size)
            run.font.name = "Cascadia Code"
            if token in keywords:
                run.font.color.rgb = KW_COLOR
            elif token.startswith('"') or token.startswith("'"):
                run.font.color.rgb = STR_COLOR
            elif token in ("Array2", "Array1", "f64", "bool", "Result", "TestResult",
                           "String", "Box", "dyn", "Send", "Sync", "Vec", "HashMap",
                           "Arc", "PyResult", "PyAny", "PyModule", "Py", "Bound",
                           "str", "np", "int", "Self",
                           "anyhow", "CITest", "Registry", "PyCITest"):
                run.font.color.rgb = STR_COLOR
            else:
                run.font.color.rgb = CODE_FG
    return tf


def add_card(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(3.1), Inches(2))
add_text(slide, "Sprint 1 Review", Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         font_size=44, bold=True)
add_text(slide, "Conditional Independence Testing", Inches(1), Inches(2.3), Inches(11), Inches(0.8),
         font_size=28, color=ACCENT)
add_text(slide, "Architecture  |  Rust CI Tests  |  Python Bindings", Inches(1), Inches(3.5),
         Inches(11), Inches(0.6), font_size=18, color=LIGHT_GRAY)
add_text(slide, "March 2026", Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         font_size=16, color=LIGHT_GRAY)

# ── SLIDE 2: Sprint Goal Recap ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, "Sprint Goal", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5))
add_card(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2))
add_text(slide,
    "Design the architecture, rewrite conditional independence tests into Rust, "
    "and create Python bindings \u2014 demonstrating that the end product is "
    "achievable and desirable (performance over the Python library).",
    Inches(1.2), Inches(1.7), Inches(10.8), Inches(1.6), font_size=20, color=LIGHT_GRAY)
add_text(slide, "Sprint Backlog Items", Inches(0.8), Inches(4.0), Inches(11), Inches(0.6),
         font_size=24, bold=True)
add_bullets(slide, [
    "\U0001f6a7  Chi-square test (in progress)",
    "\U0001f6a7  Power divergence (in progress)",
    "\u2705  Extensibility \u2014 Registry + Strategy pattern",
    "\u2705  Python bindings \u2014 PyO3 + Maturin",
    "\u2705  Git repository setup",
    "\u2705  PearsonCorrelation \u2014 fully implemented & tested",
], Inches(1.0), Inches(4.5), Inches(10), Inches(2.5), font_size=18)

# ── SLIDE 3: Architecture Overview ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, "Architecture", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5))

# Top row: three equal cards
card_w = Inches(3.7)
card_h = Inches(3.6)
gap = Inches(0.4)
x1 = Inches(0.8)
x2 = x1 + card_w + gap
x3 = x2 + card_w + gap
top_y = Inches(1.5)

# Card 1: Workspace
add_card(slide, x1, top_y, card_w, card_h)
add_text(slide, "Rust Workspace", x1 + Inches(0.2), top_y + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
         font_size=20, bold=True, color=ACCENT)
add_bullets(slide, [
    "ci-core      \u2014 Core library",
    "  \u251c strategy.rs   (CITest trait)",
    "  \u251c registry.rs   (Registry)",
    "  \u2514 ci_tests/     (Test impls)",
    "",
    "ci_python   \u2014 PyO3 bindings",
    "ci-r            \u2014 R bindings (planned)",
    "ci_js           \u2014 WASM (planned)",
], x1 + Inches(0.2), top_y + Inches(0.6), card_w - Inches(0.4), card_h - Inches(0.7),
    font_size=13, color=LIGHT_GRAY)

# Card 2: Strategy Pattern
add_card(slide, x2, top_y, card_w, Inches(0.5))
add_text(slide, "Strategy Pattern", x2 + Inches(0.2), top_y + Inches(0.05), card_w - Inches(0.4), Inches(0.4),
         font_size=20, bold=True, color=ACCENT)
add_code_block(slide, [
    "trait CITest: Send + Sync {",
    "  fn run_test(",
    "    &self,",
    "    array: Array2<f64>,",
    "    x_value: Array1<f64>,",
    "    y_value: Array1<f64>,",
    "    boolean: bool,",
    "  ) -> anyhow::Result<TestResult>;",
    "}",
], x2, top_y + Inches(0.55), card_w, card_h - Inches(0.55), font_size=12, rust_keywords=True)

# Card 3: Registry Pattern
add_card(slide, x3, top_y, card_w, card_h)
add_text(slide, "Registry Pattern", x3 + Inches(0.2), top_y + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
         font_size=20, bold=True, color=ACCENT)
add_bullets(slide, [
    "Stateless, stores boxed trait objects",
    "HashMap<String, Box<dyn CITest>>",
    "",
    "Case-insensitive lookup",
    "Easily extensible: add_to_registry()",
    "8 tests registered",
], x3 + Inches(0.2), top_y + Inches(0.6), card_w - Inches(0.4), card_h - Inches(0.7),
    font_size=13, color=LIGHT_GRAY)

# Bottom: Data Flow card spanning full width
flow_y = top_y + card_h + gap
flow_w = x3 + card_w - x1
add_card(slide, x1, flow_y, flow_w, Inches(1.6))
add_text(slide, "Data Flow", x1 + Inches(0.2), flow_y + Inches(0.1), Inches(3), Inches(0.4),
         font_size=20, bold=True, color=ACCENT)
add_text(slide,
    "Python (pandas/numpy)  \u2192  PyO3 bridge  \u2192  Rust ndarray  "
    "\u2192  CITest::run_test()  \u2192  TestResult  \u2192  Python bool / (p, coeff)",
    x1 + Inches(0.2), flow_y + Inches(0.5), flow_w - Inches(0.4), Inches(0.9),
    font_size=16, color=LIGHT_GRAY)

# ── SLIDE 4: Key Achievement ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, "Key Achievement: pearsonr in Python", Inches(0.8), Inches(0.4),
         Inches(11), Inches(0.8), font_size=36, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))
add_text(slide, "End-to-end working pipeline: Python \u2192 Rust \u2192 Python",
         Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
         font_size=22, color=GREEN, bold=True)

add_text(slide, "Usage Example", Inches(0.8), Inches(2.0), Inches(5), Inches(0.5),
         font_size=20, bold=True, color=ACCENT)
add_code_block(slide, [
    "from ci_python import PyRegistry",
    "import numpy as np",
    "",
    "registry = PyRegistry()",
    "test = registry.get_test('pearson_correlation')",
    "",
    "Z = np.column_stack([Z1, Z2])",
    "x = df['X'].to_numpy()",
    "y = df['Y'].to_numpy()",
    "",
    "result = test(Z, x, y, boolean=True)",
    "",
    "result = test(Z, x, y, boolean=False)",
], Inches(0.8), Inches(2.5), Inches(11.5), Inches(4.5), font_size=14, py_keywords=True)

# ── SLIDE 6: Performance Benchmark ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, "Performance: Rust vs Python", Inches(0.8), Inches(0.4),
         Inches(11), Inches(0.8), font_size=36, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))
add_text(slide, "PearsonCorrelation with conditioning variables (Z)",
         Inches(0.8), Inches(1.3), Inches(11), Inches(0.5), font_size=20, color=LIGHT_GRAY)

# Rust card
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.0))
add_text(slide, "Rust (via PyO3)", Inches(1.0), Inches(2.1), Inches(5), Inches(0.5),
         font_size=22, bold=True, color=GREEN)
add_text(slide, "1.37 ms", Inches(1.0), Inches(2.7), Inches(5), Inches(0.8),
         font_size=44, bold=True, color=GREEN)
add_bullets(slide, [
    "lstsq(X):          102.5 \u00b5s",
    "lstsq(Y):          149.5 \u00b5s",
    "residuals:         116.3 \u00b5s",
    "pearsonr:           84.1 \u00b5s",
    "total run_test:   452.5 \u00b5s",
], Inches(1.0), Inches(3.6), Inches(5), Inches(2.2), font_size=14, color=LIGHT_GRAY)

# Python card
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.0))
add_text(slide, "Pure Python (pgmpy)", Inches(7.0), Inches(2.1), Inches(5), Inches(0.5),
         font_size=22, bold=True, color=ORANGE)
add_text(slide, "8.06 ms", Inches(7.0), Inches(2.7), Inches(5), Inches(0.8),
         font_size=44, bold=True, color=ORANGE)
add_bullets(slide, [
    "Same test, same data, same results",
    "Identical p-values and coefficients",
    "",
    "Values match to 14 decimal places:",
    "  p = 0.20381261707...",
    "  r = 0.09023876680...",
], Inches(7.0), Inches(3.6), Inches(5), Inches(2.2), font_size=14, color=LIGHT_GRAY)

# Speedup banner
add_card(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.9), color=RGBColor(0x0A, 0x4D, 0x68))
add_text(slide, "\u26a1 ~6x faster with identical results",
         Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.6),
         font_size=24, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ── SLIDE 6: Test Coverage ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, "Test Coverage", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5))

# Rust tests card
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.5))
add_text(slide, "Rust Unit Tests (13 tests)", Inches(1.2), Inches(1.9), Inches(5.5), Inches(0.5),
         font_size=22, bold=True, color=ACCENT)
add_bullets(slide, [
    "Registry (3 tests):",
    "  \u2022 Initialization",
    "  \u2022 Lookup (case-insensitive)",
    "  \u2022 Listing all tests",
    "",
    "PearsonCorrelation (10 tests):",
    "  \u2022 Unconditional: independent & correlated",
    "  \u2022 Conditional: with confounders",
    "  \u2022 V-structure (collider) scenarios",
    "  \u2022 Multiple conditioning variables",
    "  \u2022 Both boolean and numeric output modes",
], Inches(1.2), Inches(2.5), Inches(5.3), Inches(4.5), font_size=16, color=LIGHT_GRAY)

# Python benchmark card
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5))
add_text(slide, "Python Benchmark (test.py)", Inches(7.4), Inches(1.9), Inches(5.5), Inches(0.5),
         font_size=22, bold=True, color=ACCENT)
add_bullets(slide, [
    "End-to-end validation:",
    "  \u2022 Rust vs pgmpy output comparison",
    "  \u2022 Independent variables (no conditioning)",
    "  \u2022 Conditional independence (Z1, Z2)",
    "",
    "Performance profiling:",
    "  \u2022 Per-operation timing (lstsq, residuals, pearsonr)",
    "  \u2022 Total run_test timing",
    "  \u2022 Side-by-side Rust vs Python duration",
], Inches(7.4), Inches(2.5), Inches(5.3), Inches(4.5), font_size=16, color=LIGHT_GRAY)

# ── SLIDE 7: Status & Remaining Work ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, "Status & Remaining Work", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5))

add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2))
add_text(slide, "Done", Inches(1.0), Inches(1.6), Inches(3.4), Inches(0.5),
         font_size=22, bold=True, color=GREEN)
add_bullets(slide, [
    "\u2705 Workspace architecture",
    "\u2705 CITest trait (Strategy pattern)",
    "\u2705 Registry pattern",
    "\u2705 PearsonCorrelation test",
    "\u2705 PyO3 bindings (ci_python)",
    "\u2705 Python PyRegistry wrapper",
    "\u2705 Rust & Python test suites",
    "\u2705 Git repository & workflow",
], Inches(1.0), Inches(2.2), Inches(3.4), Inches(4), font_size=14, color=LIGHT_GRAY)

add_card(slide, Inches(4.9), Inches(1.5), Inches(3.7), Inches(5.2))
add_text(slide, "Near Future", Inches(5.1), Inches(1.6), Inches(3.4), Inches(0.5),
         font_size=22, bold=True, color=ORANGE)
add_bullets(slide, [
    "\U0001f6a7 Chi-square implementation",
    "\U0001f6a7 Power divergence implementation",
    "\U0001f6a7 G-test implementation",
    "\U0001f6a7 Likelihood ratio test",
    "\U0001f6a7 Significance level config",
], Inches(5.1), Inches(2.2), Inches(3.4), Inches(4), font_size=14, color=LIGHT_GRAY)

add_card(slide, Inches(9.0), Inches(1.5), Inches(3.7), Inches(5.2))
add_text(slide, "Future", Inches(9.2), Inches(1.6), Inches(3.4), Inches(0.5),
         font_size=22, bold=True, color=ACCENT2)
add_bullets(slide, [
    "\u23f3 R bindings (extendr)",
    "\u23f3 JS/WASM bindings",
    "\u23f3 Auto test-type inference",
    "\u23f3 Additional CI tests",
], Inches(9.2), Inches(2.2), Inches(3.4), Inches(4), font_size=14, color=LIGHT_GRAY)

# Save
output = r"o:\Programming\Rust\Conditional-Independence-Testing\Sprint1_Review.pptx"
prs.save(output)
print(f"Saved: {output}")
